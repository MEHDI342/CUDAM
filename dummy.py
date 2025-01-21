from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re
from typing import Optional, Dict, Any
import logging
import sys

# Configure enterprise logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('presentation_generator.log'),
        logging.StreamHandler(sys.stdout)
    ]
)
logger = logging.getLogger(__name__)

# Enterprise configuration constants
class PresentationConfig:
    TITLE_FONT_SIZE: int = 44
    SUBTITLE_FONT_SIZE: int = 32
    BASE_CONTENT_SIZE: int = 28
    LEVEL_REDUCTION: int = 2
    SLIDE_WIDTH: int = 16
    SLIDE_HEIGHT: int = 9

    # Corporate colors
    BLACK: RGBColor = RGBColor(0, 0, 0)
    GRAY: RGBColor = RGBColor(89, 89, 89)
    ACCENT: RGBColor = RGBColor(0, 112, 192)

class EnterpriseSlideGenerator:
    """Enterprise-grade PowerPoint slide generation system with comprehensive error handling."""

    def __init__(self) -> None:
        """Initialize presentation with enterprise configuration."""
        self.prs = Presentation()
        self.prs.slide_width = Inches(PresentationConfig.SLIDE_WIDTH)
        self.prs.slide_height = Inches(PresentationConfig.SLIDE_HEIGHT)
        logger.info("Initialized EnterpriseSlideGenerator with standard configuration")

    def create_presentation(self, content: str, output_file: str) -> None:
        """
        Create complete enterprise presentation with comprehensive error handling.

        Args:
            content (str): Markdown-formatted presentation content
            output_file (str): Output file path for the presentation
        """
        try:
            slides_content = content.split('______________\n')

            for i, slide_content in enumerate(slides_content):
                if i == 0:
                    self._create_title_slide(slide_content)
                else:
                    self._create_content_slide(slide_content)

            self.prs.save(output_file)
            logger.info(f"Successfully generated presentation: {output_file}")

        except Exception as e:
            logger.error(f"Error generating presentation: {str(e)}")
            raise

    def _create_title_slide(self, content: str) -> None:
        """Create enterprise-formatted title slide with error handling."""
        try:
            lines = content.split('\n')
            title = lines[0].strip('# ')
            subtitle = lines[1].strip('## ')
            sub_subtitle = lines[2].strip('### ')

            slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1]

            title_shape.text = title
            subtitle_shape.text = f"{subtitle}\n{sub_subtitle}"

            self._format_title_text(title_shape.text_frame)
            self._format_subtitle_text(subtitle_shape.text_frame)

        except Exception as e:
            logger.error(f"Error creating title slide: {str(e)}")
            raise

    def _create_content_slide(self, content: str) -> None:
        """Create enterprise-formatted content slide with comprehensive error handling."""
        try:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
            shapes = slide.shapes

            sections = content.split('\n\n')
            if not sections:
                raise ValueError("Invalid slide content: No sections found")

            title = sections[0].strip('# ')
            title_shape = shapes.title
            title_shape.text = title
            self._format_title_text(title_shape.text_frame)

            body_shape = shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()

            for section in sections[1:]:
                self._add_formatted_content(tf, section)

        except Exception as e:
            logger.error(f"Error creating content slide: {str(e)}")
            raise

    def _format_title_text(self, frame) -> None:
        """Apply enterprise title formatting with error handling."""
        try:
            p = frame.paragraphs[0]
            p.font.size = Pt(PresentationConfig.TITLE_FONT_SIZE)
            p.font.bold = True
            p.font.color.rgb = PresentationConfig.BLACK
        except Exception as e:
            logger.error(f"Error formatting title text: {str(e)}")
            raise

    def _format_subtitle_text(self, frame) -> None:
        """Apply enterprise subtitle formatting with error handling."""
        try:
            p = frame.paragraphs[0]
            p.font.size = Pt(PresentationConfig.SUBTITLE_FONT_SIZE)
            p.font.color.rgb = PresentationConfig.GRAY
        except Exception as e:
            logger.error(f"Error formatting subtitle text: {str(e)}")
            raise

    def _add_formatted_content(self, tf, content: str) -> None:
        """Add formatted content with proper hierarchy and error handling."""
        try:
            lines = content.strip().split('\n')

            for line in lines:
                if not line.strip():
                    continue

                level = self._determine_level(line)
                clean_line = self._clean_line(line)

                p = tf.add_paragraph()
                p.text = clean_line
                p.level = level

                self._format_paragraph(p, level)

        except Exception as e:
            logger.error(f"Error adding formatted content: {str(e)}")
            raise

    def _determine_level(self, line: str) -> int:
        """Determine content hierarchy level with validation."""
        try:
            if line.startswith('o'):
                return 2
            if line.startswith('•'):
                return 1
            return 0
        except Exception as e:
            logger.error(f"Error determining content level: {str(e)}")
            raise

    def _clean_line(self, line: str) -> str:
        """Clean line formatting markers with validation."""
        try:
            return line.lstrip('o• ').strip()
        except Exception as e:
            logger.error(f"Error cleaning line: {str(e)}")
            raise

    def _format_paragraph(self, p, level: int) -> None:
        """Apply enterprise paragraph formatting with error handling."""
        try:
            p.font.size = Pt(PresentationConfig.BASE_CONTENT_SIZE - (level * PresentationConfig.LEVEL_REDUCTION))
            p.font.color.rgb = PresentationConfig.BLACK
            if level == 0:
                p.font.bold = True
        except Exception as e:
            logger.error(f"Error formatting paragraph: {str(e)}")
            raise

def create_presentation(content: str, output_file: str) -> None:
    """
    Create enterprise presentation with comprehensive error handling.

    Args:
        content (str): Markdown-formatted presentation content
        output_file (str): Output file path for the presentation
    """
    try:
        generator = EnterpriseSlideGenerator()
        generator.create_presentation(content, output_file)
    except Exception as e:
        logger.error(f"Error in presentation creation: {str(e)}")
        raise

def main():
    """Main execution with complete content implementation."""
    markdown_content = """# L'Intelligence Artificielle : Solutions et Applications
## Analyse d'Implémentation 
### SOCOBAT - Transformation Digitale 2024

______________

# Domaines d’application de l’IA : Exemples spécifiques
## Étude de Marché sur l’Utilisation de l’IA dans l’Assistance Client

### Exemples Spécifiques

• Exemple 1 : Société Générale
o Utilisation de chatbots IA pour répondre aux demandes des clients en ligne.
o Résultat : réduction significative des délais de réponse, amélioration de la satisfaction client grâce à des réponses rapides et précises.

• Exemple 2 : Amazon
o Intégration d’assistants virtuels pour gérer les interactions client.
o Résultat : personnalisation accrue des échanges et réduction des coûts opérationnels en automatisant les requêtes les plus courantes.

______________

# Production
## Technologies Avancées

### Exemples de Production

• Exemple 1 : Bouygues Construction
o Utilisation de l’IA pour la gestion des chantiers, notamment pour surveiller les délais et optimiser les ressources.
o Résultat : réduction des retards dans les projets et meilleure coordination des équipes.

• Exemple 2 : VESTACK
o Intégration d’algorithmes prédictifs pour optimiser les processus de construction écologique.
o Résultat : diminution des coûts énergétiques et des émissions de CO₂ grâce à des simulations en temps réel.

______________

# Finance et Comptabilité
## Optimisation des Processus

### Exemples en Finance

• Exemple 1 : Cogedis
o Automatisation des processus de comptabilité grâce à l’IA.
o Résultat : amélioration de 30 % de l’efficacité des opérations et réduction des erreurs dans les calculs financiers.

• Exemple 2 : ENGIE
o Analyse prédictive des flux financiers pour anticiper les tendances du marché.
o Résultat : prise de décision plus rapide et fiable grâce à une meilleure visibilité des données financières.

______________

# Gestion d’approvisionnement
## Solutions Logistiques

### Exemple en Gestion d’Approvisionnement

• Exemple unique : Suffolk Construction
o Mise en place d’algorithmes d’IA pour prévoir les besoins en matériaux et optimiser les délais de livraison.
o Résultat : réduction des coûts logistiques et anticipation des pénuries grâce à des données en temps réel.

______________

# Sécurité des employés
## Surveillance Intelligente

### Exemples en Sécurité

• Exemple 1 : IBM
o Déploiement de systèmes IA pour surveiller les comportements et détecter les risques sur les lieux de travail.
o Résultat : réduction des incidents de sécurité et meilleure conformité aux normes.

• Exemple 2 : Siemens
o Utilisation de l’IA pour analyser les données des équipements et prévenir les défaillances.
o Résultat : diminution des accidents liés à l’utilisation de machines défectueuses.

______________

# Startups dans la région de Socobat : Présentation et témoignages
## Écosystème Local

### Finance et Comptabilité : Yooz (Lyon)

• Présentation :
Yooz se concentre sur l’automatisation comptable, simplifiant la gestion des factures.
o Témoignage :
Brico Privé a constaté une transparence accrue dans ses processus grâce à Yooz, améliorant les relations avec les fournisseurs.

### Logistique et Gestion de la Chaîne d'Approvisionnement

• Kardinal (Lyon) :
Optimisation des tournées de livraison à l’aide d’algorithmes avancés.
o Témoignage :
DPD France a observé une augmentation de 4 % dans le respect des créneaux horaires grâce aux solutions Kardinal.

• Aitenders (Saint-Étienne) :
Automatisation des appels d’offres pour réduire les tâches manuelles et améliorer les sélections.
o Témoignage :
Une entreprise cliente a réduit de 90 % les opportunités non pertinentes, se concentrant sur les offres qualifiées.

### Service Client, Sécurité et Production : Wevioo (Lyon)

Wevioo accompagne les entreprises dans leur transformation digitale, en optimisant la production, la maintenance prédictive et la sécurité des employés grâce à des solutions basées sur l’intelligence artificielle.
o Témoignage :
Une entreprise industrielle a collaboré avec Wevioo pour intégrer des systèmes de maintenance prédictive, permettant de réduire les coûts et d’améliorer la productivité globale.

______________

# L’impact de l’intégration de l’IA : Pourcentages clés
## Métriques Clés

### Finance et Comptabilité

• Réduction des coûts administratifs : 20-30 %.
• Amélioration de la productivité : 25 %.

### Logistique et Gestion de la chaîne d'approvisionnement

• Réduction des coûts logistiques : 15-20 %.
• Amélioration de l’efficacité des chaînes d’approvisionnement : 20 %.

### Service Client et Assistance

• Réduction des coûts de service client : 30 %.
• Temps de réponse amélioré : 50 %.

### Sécurité des Employés

• Réduction des incidents liés à la sécurité : 10-15 %.
• Amélioration de la conformité : 20 %.

### Production

• Réduction des coûts de maintenance : 15-30 %.
• Amélioration de la productivité : 20 %.

______________

# Défis de l’intégration de l’IA dans l’entreprise et solutions
## Obstacles et Solutions

### Défi 1 : Coût élevé de mise en œuvre

• Solution :
o Déployer des projets pilotes avant un déploiement global.
o Rechercher des subventions et aides publiques.

### Défi 2 : Résistance au changement interne

• Solution :
o Impliquer les employés dès le début.
o Proposer des formations adaptées.

### Défi 3 : Données insuffisantes ou de mauvaise qualité

• Solution :
o Mettre en place un processus de collecte et de nettoyage des données.
o Centraliser les données pour une meilleure accessibilité.

### Défi 4 : Complexité technologique

• Solution :
o Collaborer avec des experts en IA.
o Simplifier les processus pour faciliter l’adoption.

### Défi 5 : Risques éthiques et réglementaires

• Solution :
o Garantir la transparence des algorithmes.
o Respecter les réglementations, comme le RGPD.

### Défi 6 : Temps nécessaire pour constater des résultats

• Solution :
o Fixer des attentes réalistes et des objectifs intermédiaires.
o Mesurer les résultats par des indicateurs de performance.

### Défi 7 : Cybersécurité

• Solution :
o Investir dans des outils de sécurité avancés.
o Former les équipes aux bonnes pratiques 
"""

    try:
        create_presentation(markdown_content, 'SOCOBAT_IA_Implementation.pptx')
        logger.info("Presentation generation completed successfully")
    except Exception as e:
        logger.error(f"Fatal error in presentation generation: {str(e)}")
        raise

if __name__ == "__main__":
    main()
